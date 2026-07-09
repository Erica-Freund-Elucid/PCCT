"""Generate the PCCT Gate-3 statistical-method report (Word) + summary deck (PPT).

Computes the variance-components decomposition on the canonical (vessel-overlap)
paired data for v1 (original) and v2 (2026-07-07), both scales for plaque, and
writes:
    PCCT_Statistical_Methods_and_Results.docx
    PCCT_Statistical_Methods_and_Results.pptx   (gitignored; regenerate from here)

Run from repo root:  python scripts/generate_stat_report.py
"""
import csv, math, os
import numpy as np

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), os.pardir))
V1C = os.path.join(ROOT, "gate_results_v1_original", "paired_data.csv")
V1S = os.path.join(ROOT, "gate_results_v1_original", "paired_data_subsegment.csv")
V2C = os.path.join(ROOT, "gate_results", "paired_data.csv")
V2S = os.path.join(ROOT, "gate_results", "paired_data_subsegment.csv")
NBOOT, SEED = 2000, 42

# OQ inter-observer references (730-CVV-040 / ATTACHMENT 2 patient level): (pt, lo, hi)
OQ = {
    "LumenVol":        {"lab": "Lumen",       "log": (7.5, 5.9, 9.3),   "ut": (9.32, 7.09, 11.58)},
    "WallVol":         {"lab": "Wall",        "log": (13.1, 10.2, 16.0),"ut": (23.50, 17.88, 29.71)},
    "VesselVol":       {"lab": "Vessel",      "log": (8.7, 6.8, 10.6),  "ut": (10.18, 7.83, 12.38)},
    "CALCVol":         {"lab": "CALC",        "log": (4.5, 3.5, 5.4),   "ut": (25.78, 18.70, 35.75), "bias_ci": (-0.01, 0.05)},
    "LRNCVol":         {"lab": "LRNC",        "log": (5.4, 4.1, 6.9),   "ut": (78.80, 48.09, 163.40), "bias_ci": (-0.06, -0.01)},
    "NonCALCMATXVol":  {"lab": "NonCALC Matrix","log": (13.6, 10.8, 16.4),"ut": (32.61, 24.87, 41.94), "bias_ci": (-0.19, 0.05)},
    "TotalPlaqueVolume": {"lab": "Total Plaque","log": (13.2, 10.4, 16.0),"ut": (27.08, 20.68, 34.74), "bias_ci": (-0.19, 0.08)},
}
PROCESS = ["LumenVol", "WallVol", "VesselVol"]
PLAQUE = ["CALCVol", "LRNCVol", "NonCALCMATXVol", "TotalPlaqueVolume"]


def _f(r, k):
    try:
        return float(r[k])
    except Exception:
        return None


def norm_pairs(rows, var):
    pv, ev = [], []
    for r in rows:
        p, e = _f(r, f"PCCT_{var}"), _f(r, f"EID_{var}")
        pl, el = _f(r, "PCCT_Len"), _f(r, "EID_Len")
        if None in (p, e, pl, el) or pl <= 0 or el <= 0:
            continue
        pv.append(p / pl); ev.append(e / el)
    return np.array(pv), np.array(ev)


def raw_pairs(rows, var):
    """Raw (non-length-normalized) totals, same patient filter as norm_pairs."""
    pv, ev = [], []
    for r in rows:
        p, e = _f(r, f"PCCT_{var}"), _f(r, f"EID_{var}")
        pl, el = _f(r, "PCCT_Len"), _f(r, "EID_Len")
        if None in (p, e, pl, el) or pl <= 0 or el <= 0:
            continue
        pv.append(p); ev.append(e)
    return pv, ev


def overlap(a, b):
    return a[0] <= b[1] and b[0] <= a[1]


def compute(path):
    rows = list(csv.DictReader(open(path, encoding="utf-8")))
    rng = np.random.RandomState(SEED)
    out = {"n": None}
    for var, meta in OQ.items():
        pv, ev = norm_pairs(rows, var)
        pvraw, evraw = raw_pairs(rows, var)
        out["n"] = len(pv)
        d = pv - ev
        dl = np.log(pv + 1) - np.log(ev + 1)
        n = len(d)
        # log-scale terms
        s2_oq = math.log(1 + (meta["log"][0] / 100) ** 2)
        rec = {
            "delta": float(dl.mean()),
            "s2_total": float(np.mean(dl ** 2) / 2),
            "s2_rand": float(np.var(dl, ddof=1) / 2),
            "s2_oq": s2_oq,
        }
        rec["s2_scan"] = rec["s2_rand"] - s2_oq
        wl = lambda s2: math.sqrt(math.exp(max(s2, 0)) - 1) * 100
        rec["wcv_total_log"] = wl(rec["s2_total"])
        rec["wcv_rand_log"] = wl(rec["s2_rand"])
        rec["wcv_scan_log"] = wl(rec["s2_scan"])
        # bootstrap CIs (scanner-term wcv, and sigma2_scan)
        bt_rand, bt_scan = [], []
        for _ in range(NBOOT):
            idx = rng.randint(0, n, n)
            db = dl[idx]
            s2r = float(np.var(db, ddof=1) / 2)
            bt_rand.append(wl(s2r)); bt_scan.append(s2r - s2_oq)
        rec["wcv_rand_log_ci"] = (float(np.percentile(bt_rand, 2.5)), float(np.percentile(bt_rand, 97.5)))
        rec["s2_scan_ci"] = (float(np.percentile(bt_scan, 2.5)), float(np.percentile(bt_scan, 97.5)))
        rec["wcv_scan_ci"] = (wl(rec["s2_scan_ci"][0]), wl(rec["s2_scan_ci"][1]))
        # Gate 4 systematic bias (log-scale BA, length-normalized) vs OQ Table 6
        bt_bias = [float(np.mean(dl[rng.randint(0, n, n)])) for _ in range(NBOOT)]
        rec["bias_ci"] = (float(np.percentile(bt_bias, 2.5)), float(np.percentile(bt_bias, 97.5)))
        oqb = meta.get("bias_ci")
        rec["bias_overlap"] = overlap(rec["bias_ci"], oqb) if oqb else None
        rec["gate3_pass_log"] = overlap(rec["wcv_rand_log_ci"], meta["log"][1:])
        rec["scan_a_pass"] = rec["wcv_scan_log"] <= meta["log"][0]     # (a) scanner wCV <= OQ
        rec["scan_b_pass"] = rec["s2_scan_ci"][0] <= 0                 # (b) CI includes 0
        # untransformed (for plaque both-scales), length-normalized
        m = float(np.mean(np.concatenate([pv, ev])))
        s2r_ut = float(np.var(d, ddof=1) / 2)
        rec["mean_ut"] = m
        rec["wcv_rand_ut"] = math.sqrt(s2r_ut) / m * 100 if m else 0
        # untransformed systematic bias (Gate 4, non-log). Report on RAW (non-
        # normalized) volumes to match gate_summary/HTML and the traditional Gate 4.
        praw = np.array([p for p, e in zip(pvraw, evraw)])
        eraw = np.array([e for p, e in zip(pvraw, evraw)])
        draw = praw - eraw
        mraw = float(np.mean(np.concatenate([praw, eraw]))) if len(praw) else 0.0
        rec["ut_bias_raw"] = float(draw.mean()) if len(draw) else 0.0
        rec["ut_bias_pct"] = abs(rec["ut_bias_raw"]) / mraw * 100 if mraw else 0.0
        rec["ut_bias_pass"] = rec["ut_bias_pct"] < 10.0   # plaque project threshold
        bt_ut = []
        for _ in range(NBOOT):
            idx = rng.randint(0, n, n)
            dd = d[idx]; mm = float(np.mean(np.concatenate([pv[idx], ev[idx]])))
            bt_ut.append(math.sqrt(np.var(dd, ddof=1) / 2) / mm * 100 if mm else 0)
        rec["wcv_rand_ut_ci"] = (float(np.percentile(bt_ut, 2.5)), float(np.percentile(bt_ut, 97.5)))
        rec["gate3_pass_ut"] = overlap(rec["wcv_rand_ut_ci"], meta["ut"][1:])
        out[var] = rec
    return out


R = {"v1c": compute(V1C), "v1s": compute(V1S), "v2c": compute(V2C), "v2s": compute(V2S)}
R1, R2 = R["v1c"], R["v2c"]  # canonical (back-compat)
COLS = [("v1c", "v1 canonical"), ("v1s", "v1 sub-seg"), ("v2c", "v2 canonical"), ("v2s", "v2 sub-seg")]


# ─────────────────────────────────────────────────────────── Word report ──
def build_docx():
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = Document()
    st = doc.styles["Normal"]; st.font.name = "Calibri"; st.font.size = Pt(10.5)

    def h(txt, lvl=1):
        doc.add_heading(txt, level=lvl)

    def p(txt, italic=False, bold=False):
        par = doc.add_paragraph(); run = par.add_run(txt)
        run.italic = italic; run.bold = bold
        return par

    def table(headers, rows, style="Light Grid Accent 1"):
        t = doc.add_table(rows=1, cols=len(headers)); t.style = style
        for i, hh in enumerate(headers):
            c = t.rows[0].cells[i]; c.text = ""
            r = c.paragraphs[0].add_run(hh); r.bold = True; r.font.size = Pt(9)
        for row in rows:
            cells = t.add_row().cells
            for i, v in enumerate(row):
                cells[i].text = ""
                rr = cells[i].paragraphs[0].add_run(str(v)); rr.font.size = Pt(9)
        return t

    title = doc.add_heading("PCCT Scanner Qualification — Statistical Method & Results", level=0)
    p("Gate 3 (quantitative reproducibility) — cross-scanner PCCT vs EID vs the B.1P Delta "
      "Validation OQ (730-CVV-040). Canonical vessel-overlap region, length-normalized, "
      f"N = {R2['n']} paired patients. Two data vintages: v1 (original) and v2 (2026-07-07 re-work). "
      "Preliminary (target N ≥ 30).", italic=True)

    h("1. Purpose & scope")
    p("Establish, on a like-for-like basis with the validated B.1P inter-operator reproducibility, "
      "whether cross-scanner (PCCT vs EID) quantitative reproducibility is non-inferior to the OQ "
      "limits. The systematic modality difference (bias) is separated from random dispersion and "
      "assessed in Gate 4; this document covers the Gate-3 reproducibility (wCV) method and the "
      "scanner-attributable variance decomposition.")

    h("2. Measurement model")
    p("For patient i, scanner s ∈ {PCCT, EID}, reader r, on the log(x+1), length-normalized scale:")
    p("    Y_isr = μ + P_i + S_s + (PS)_is + R_r + e_isr", bold=True)
    table(["Term", "Meaning", "Role"], [
        ["P_i", "patient (between-subject) effect", "not part of within-subject reproducibility"],
        ["S_s", "scanner main effect (systematic modality bias)", "FIXED → Gate 4 (not variance)"],
        ["(PS)_is", "patient × scanner interaction", "σ²_scanner — the scanner-attributable variance"],
        ["R_r", "reader (inter-operator) effect", "σ²_R — measured by the OQ"],
        ["e_isr", "residual / repeat error", "σ²_e — measured by the OQ"],
    ])
    p("Only one read per patient·scanner is available, so only the paired difference "
      "d_i = Y_PCCT − Y_EID is observed. P_i cancels, giving "
      "d_i = Δ + [(PS) difference] + [R − R′] + [e − e′], hence Var(d)/2 = σ²_PS + σ²_R + σ²_e.")

    h("3. Estimators & acceptance")
    p("wCV (within-subject CV) — variance-component / random-effects estimator per Quan & Shih "
      "(1996), matching the OQ (730-CVV-040 §11: log(x+1) + linear mixed model). For paired (k=2) "
      "data the within-subject variance component is σ²_w = mean(d²)/2.")
    table(["Quantity", "Formula", "Purpose"], [
        ["Systematic bias Δ", "mean(d) (log scale)", "Gate 4 (BA bias)"],
        ["σ²_total", "mean(d²)/2 = σ²_random + Δ²/2", "within-subject var incl. bias"],
        ["σ²_random (scanner term)", "Var(d)/2", "bias-removed dispersion (Gate 3 basis)"],
        ["σ²_OQ,within", "ln(1 + (OQ_logCV/100)²)", "reader+repeat, imported from OQ"],
        ["σ²_scanner", "Var(d)/2 − σ²_OQ,within", "scanner×patient interaction"],
        ["wCV (log)", "√(exp(σ²)−1)·100", "report/compare on log scale"],
        ["wCV (untransf.)", "√σ²_w / grand-mean · 100", "report on untransformed scale"],
    ])
    p("95% CIs: 2000-sample bootstrap over patients (seed 42). Acceptance — Gate 3: PCCT wCV 95% CI "
      "OVERLAPS the delta-OQ wCV CI (scanner-term basis). Scanner-attributable: (a) scanner wCV ≤ OQ "
      "inter-observer limit, and (b) σ²_scanner 95% CI includes 0 (not distinguishable from zero). "
      "The scanner-attributable subtraction is done on the log scale (there the CV↔variance relation "
      "is mean-independent, so the OQ term is transportable across studies).")

    h("4. OQ reference (730-CVV-040 / ATTACHMENT 2, patient level)")
    table(["Endpoint", "OQ log-wCV [95% CI]", "OQ untransf-wCV [95% CI]"],
          [[OQ[v]["lab"], f'{OQ[v]["log"][0]}% [{OQ[v]["log"][1]}, {OQ[v]["log"][2]}]',
            f'{OQ[v]["ut"][0]}% [{OQ[v]["ut"][1]}, {OQ[v]["ut"][2]}]'] for v in OQ])

    h("5. Results — variance decomposition (log scale), canonical & sub-segment")
    def dec_rows(rr):
        out = []
        for v in OQ:
            r = rr[v]
            out.append([OQ[v]["lab"], f'{r["delta"]:+.3f}', f'{r["s2_total"]:.5f}', f'{r["s2_rand"]:.5f}',
                        f'{r["s2_oq"]:.5f}', f'{r["s2_scan"]:.5f}', f'{r["wcv_total_log"]:.1f}',
                        f'{r["wcv_rand_log"]:.1f}', f'{r["wcv_scan_log"]:.1f}'])
        return out
    hdr = ["Endpoint", "Δ bias", "σ²_total", "σ²_rand", "σ²_OQ", "σ²_scan",
           "wCV_tot", "wCV_scanT", "wCV_attrib"]
    for key, lab in COLS:
        p(f"{lab}  (N={R[key]['n']}):", bold=True); table(hdr, dec_rows(R[key]))
    p("Key: within a vintage, σ²_random / σ²_scanner shrink from canonical → sub-segment (the "
      "traced-extent term is removed). Across vintages (canonical), σ²_random is stable while Δ "
      "(systematic bias) grew v1→v2 — the 07-07 re-work changed the bias (Gate 4), not the random "
      "reproducibility (Gate 3).", italic=True)

    def cell_g3(r, scale):
        if scale == "log":
            return f'{r["wcv_rand_log"]:.1f} {"P" if r["gate3_pass_log"] else "F"}'
        return f'{r["wcv_rand_ut"]:.1f} {"P" if r["gate3_pass_ut"] else "F"}'

    h("6. Gate 3 acceptance — scanner-term wCV vs OQ (95% CI overlap; P=pass/F=fail)")
    hdr3 = ["Endpoint", "OQ wCV"] + [c[1] for c in COLS]
    p("Log scale — process outputs (primary) + plaque:", bold=True)
    rows = []
    for v in PROCESS + PLAQUE:
        rows.append([OQ[v]["lab"], f'{OQ[v]["log"][0]} [{OQ[v]["log"][1]},{OQ[v]["log"][2]}]']
                    + [cell_g3(R[k][v], "log") for k, _ in COLS])
    table(hdr3, rows)
    p("Untransformed scale — plaque (reported on both scales):", bold=True)
    rows = []
    for v in PLAQUE:
        rows.append([OQ[v]["lab"], f'{OQ[v]["ut"][0]} [{OQ[v]["ut"][1]},{OQ[v]["ut"][2]}]']
                    + [cell_g3(R[k][v], "ut") for k, _ in COLS])
    table(hdr3, rows)
    p("Cells: scanner-term wCV% and overlap verdict; 95% CIs are in §5 / gate3_comparison.csv.",
      italic=True)

    h("7. Scanner-attributable variance vs OQ — canonical vs sub-segment (a: wCV≤OQ, b: CI∋0)")
    hdrsa = ["Endpoint", "OQ limit"] + [c[1] for c in COLS]
    rows = []
    for v in OQ:
        cells = []
        for k, _ in COLS:
            r = R[k][v]
            cells.append(f'{r["wcv_scan_log"]:.1f} {"P" if r["scan_a_pass"] else "F"}/{"P" if r["scan_b_pass"] else "F"}')
        rows.append([OQ[v]["lab"], f'{OQ[v]["log"][0]}'] + cells)
    table(hdrsa, rows)
    p("Cells: scanner-attributable wCV%  a/b (P/F). Canonical still contains the traced-extent "
      "term; sub-segment removes it — scanner-attributable variance drops accordingly and is "
      "≤ OQ / CI-includes-0 for nearly all endpoints on sub-segment.", italic=True)

    h("8. Gate 4 — systematic bias, both scales (log Bland-Altman & untransformed)")
    p("8a. LOG scale — systematic PCCT−EID bias Δ on the log(x+1), length-normalized scale. "
      "Acceptance: PCCT bias 95% CI overlaps the 730-CVV-040 Table 6 inter-observer bias 95% CI "
      "(which essentially includes 0). Plaque only (Table 6 = plaque BA refs). P=overlap/pass.",
      bold=True)
    hdrb = ["Endpoint", "OQ bias [95% CI]"] + [c[1] for c in COLS]
    rows = []
    for v in PLAQUE:
        oqb = OQ[v]["bias_ci"]
        cells = []
        for k, _ in COLS:
            r = R[k][v]
            cells.append(f'{r["delta"]:+.3f} [{r["bias_ci"][0]:+.3f},{r["bias_ci"][1]:+.3f}] '
                         f'{"P" if r["bias_overlap"] else "F"}')
        rows.append([OQ[v]["lab"], f'[{oqb[0]}, {oqb[1]}]'] + cells)
    table(hdrb, rows)
    p("8b. UNTRANSFORMED (non-log) scale — bias as % of mean on RAW (non-normalized) volumes "
      "(matches gate_summary / HTML report / the traditional Gate 4), vs the project |bias| < 10% "
      "of mean threshold. P=pass. No OQ untransformed BA-bias reference — project-specific criterion. "
      "e.g. NonCALC Matrix v2 canonical bias ≈ −108.9 mm³ (47% of mean).", bold=True)
    hdru = ["Endpoint", "threshold"] + [c[1] for c in COLS]
    rows = []
    for v in PLAQUE:
        cells = [f'{R[k][v]["ut_bias_pct"]:.1f}% {"P" if R[k][v]["ut_bias_pass"] else "F"}' for k, _ in COLS]
        rows.append([OQ[v]["lab"], "<10%"] + cells)
    table(hdru, rows)
    p("Full BA bias/LoA/proportional-bias detail and the OQ-overlay plots are in "
      "gate_results/qualification_report.html (canonical) and the gate_summary.txt sub-segment "
      "section + gate_results/bland_altman_plots_subsegment/ (sub-segment). Note the systematic "
      "plaque bias grew v1→v2 (07-07 re-work) and does NOT shrink to overlap on sub-segment — the "
      "bias is within the shared centerline, not at the traced-extent tail (consistent with the "
      "earlier sub-segment finding).", italic=True)

    h("9. Interpretation")
    for b in [
        "With the scanner term (systematic bias removed), the process outputs (Lumen, Wall, Vessel) "
        "pass the OQ CI-overlap on BOTH data vintages, canonical region — no sub-segment alignment "
        "required. The essential ingredient is separating the systematic scanner bias into Gate 4.",
        "The variance decomposition shows σ²_random / σ²_scanner are vintage-stable; the 07-07 "
        "re-work increased the systematic bias Δ, not the random dispersion.",
        "Plaque reproducibility depends on scale: on the untransformed scale (wide OQ CIs) all "
        "plaque endpoints overlap; on the log scale the tightest-limit component (CALC) is borderline "
        "and just fails on v2. Both scales are reported.",
        "Scanner-attributable variance (Var(d)/2 − σ²_OQ) is ≤ the OQ limit / CI-includes-0 for most "
        "endpoints on the canonical region, and on the extent-controlled sub-segment region "
        "(regenerated on v2, N=24) it is ≤ OQ / CI-includes-0 for nearly all endpoints — i.e. after "
        "removing reader+repeat (OQ) and traced extent (sub-segment), the residual scanner-attributable "
        "variability is within the OQ envelope.",
    ]:
        doc.add_paragraph(b, style="List Bullet")

    h("10. Caveats & future work")
    for b in [
        "N = 25 canonical / 24 sub-segment, preliminary (target ≥ 30); passes are CI-overlap-driven "
        "with wide CIs.",
        "Canonical region retains a traced-extent differential (PCCT traces ~+26% longer within "
        "shared vessels; Gate 2 REVIEW). The sub-segment (distance-from-ostium) intersection removes "
        "it; regenerated on v2 (2026-07-07) data, N=24 (PT-124 no vessel overlap, PT-136 duplicate "
        "LeftCoronary in workitem.json — both excluded). Pipeline: scripts/subsegment/.",
        "Scanner-attributable subtraction assumes the OQ reader+repeat variance is transportable to "
        "the PCCT setting, and that PCCT/EID reads are by independent readers (single read per "
        "patient·scanner; true scanner isolation would need replicate reads).",
    ]:
        doc.add_paragraph(b, style="List Bullet")

    h("References")
    for r in ["730-CVV-040 v0.1 — B.1P Delta Validation OQ Report (+ ATTACHMENT 2, patient-level results).",
              "4-B1P-033 v2.0 — B.1P Operational Qualification Report.",
              "Quan H, Shih WJ (1996). Assessing reproducibility by the within-subject coefficient of "
              "variation with random effects models. Biometrics 52:1195-1203.",
              "Shoukri MM, Elkum N, Walter SD (2006). Interval estimation and optimal design for the "
              "within-subject coefficient of variation. BMC Med Res Methodol 6:24."]:
        doc.add_paragraph(r, style="List Bullet")

    out = os.path.join(ROOT, "PCCT_Statistical_Methods_and_Results.docx")
    doc.save(out); print("Wrote", out)


# ─────────────────────────────────────────────────────────────── PPT deck ──
def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    NAVY = RGBColor(0x1f, 0x29, 0x37); BLUE = RGBColor(0x25, 0x63, 0xeb)
    GREEN = RGBColor(0x16, 0xa3, 0x4a); RED = RGBColor(0xdc, 0x26, 0x26)

    def slide(title):
        s = prs.slides.add_slide(BLANK)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = title
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
        return s

    def bullets(s, items, top=1.3, size=18, left=0.7, width=12.0):
        tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.6))
        tf = tb.text_frame; tf.word_wrap = True
        for i, (txt, lvl, color) in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.level = lvl
            run = para.add_run(); run.text = ("• " if lvl == 0 else "– ") + txt
            run.font.size = Pt(size - 2 * lvl);
            if color: run.font.color.rgb = color
        return s

    def table_slide(title, headers, rows, colw=None, size=12):
        s = slide(title)
        nr, nc = len(rows) + 1, len(headers)
        gt = s.shapes.add_table(nr, nc, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4 * nr)).table
        if colw:
            for i, w in enumerate(colw): gt.columns[i].width = Inches(w)
        for i, hh in enumerate(headers):
            c = gt.cell(0, i); c.text = hh
            c.text_frame.paragraphs[0].runs[0].font.size = Pt(size); c.text_frame.paragraphs[0].runs[0].font.bold = True
        for ri, row in enumerate(rows, 1):
            for ci, v in enumerate(row):
                c = gt.cell(ri, ci); c.text = str(v)
                rn = c.text_frame.paragraphs[0].runs[0]; rn.font.size = Pt(size)
                if v == "PASS": rn.font.color.rgb = GREEN; rn.font.bold = True
                if v == "FAIL": rn.font.color.rgb = RED; rn.font.bold = True
        return s

    # 1 title
    s = slide("PCCT Scanner Qualification — Statistical Method & Results")
    bullets(s, [
        ("Gate 3 quantitative reproducibility: PCCT vs EID vs B.1P Delta OQ (730-CVV-040)", 0, None),
        (f"Canonical vessel-overlap, length-normalized, N = {R2['n']} paired (preliminary)", 0, None),
        ("Two data vintages: v1 (original) and v2 (2026-07-07 re-work)", 0, None),
        ("Variance-component estimator (Quan & Shih 1996) + scanner term", 0, BLUE),
    ], top=2.2, size=20)

    # 2 model
    s = slide("Measurement model")
    bullets(s, [
        ("Y_isr = μ + P_i + S_s + (PS)_is + R_r + e_isr    (log(x+1), length-normalized)", 0, BLUE),
        ("P_i patient (between-subject) — not within-subject reproducibility", 0, None),
        ("S_s scanner main effect = systematic modality bias → FIXED, assessed in Gate 4", 0, None),
        ("(PS)_is patient×scanner interaction = σ²_scanner (scanner-attributable)", 0, None),
        ("R_r reader + e residual = σ²_OQ,within (measured by the OQ inter-observer study)", 0, None),
        ("One read per patient·scanner → observe d = Y_PCCT − Y_EID;  Var(d)/2 = σ²_scanner + σ²_OQ", 0, NAVY),
    ], top=1.5, size=18)

    # 3 estimators
    s = slide("Estimators & acceptance")
    bullets(s, [
        ("wCV = within-subject CV, variance-component:  σ²_w = mean(d²)/2", 0, None),
        ("log wCV = √(exp(σ²)−1)·100 ;  untransformed wCV = √σ²_w / mean · 100", 1, None),
        ("Scanner term: σ²_random = Var(d)/2 (systematic bias Δ removed; Δ → Gate 4)", 0, None),
        ("σ²_total = σ²_random + Δ²/2", 1, None),
        ("Scanner-attributable: σ²_scanner = Var(d_log)/2 − σ²_OQ,within", 0, None),
        ("σ²_OQ,within = ln(1 + (OQ_logCV/100)²)  from 730-CVV-040 ATTACHMENT 2", 1, None),
        ("Acceptance — Gate 3: PCCT wCV 95% CI overlaps delta-OQ (scanner-term basis)", 0, GREEN),
        ("Scanner-attributable: (a) scanner wCV ≤ OQ  and  (b) σ²_scanner 95% CI includes 0", 0, GREEN),
        ("95% CIs: 2000-sample bootstrap over patients", 1, None),
    ], top=1.4, size=16)

    # 4 gate3 log results v2
    def g3(R, scale):
        vs = PROCESS + PLAQUE if scale == "log" else PLAQUE
        rows = []
        for v in vs:
            r = R[v]
            if scale == "log":
                wcv, ci, ok, oq = r["wcv_rand_log"], r["wcv_rand_log_ci"], r["gate3_pass_log"], OQ[v]["log"]
            else:
                wcv, ci, ok, oq = r["wcv_rand_ut"], r["wcv_rand_ut_ci"], r["gate3_pass_ut"], OQ[v]["ut"]
            rows.append([OQ[v]["lab"], f'{oq[0]} [{oq[1]}, {oq[2]}]',
                         f'{wcv:.1f} [{ci[0]:.1f}, {ci[1]:.1f}]', "PASS" if ok else "FAIL"])
        return rows
    table_slide("Gate 3 — scanner-term wCV vs OQ (log scale, v2 / 2026-07-07)",
                ["Endpoint", "OQ wCV [95% CI]", "PCCT wCV [95% CI]", "Overlap"], g3(R2, "log"),
                colw=[3.0, 3.6, 3.9, 1.8])
    table_slide("Plaque reported on both scales (v2) — untransformed",
                ["Endpoint", "OQ wCV [95% CI]", "PCCT wCV [95% CI]", "Overlap"], g3(R2, "ut"),
                colw=[3.0, 3.9, 3.9, 1.5])

    # 5 decomposition insight
    dec = []
    for v in ["WallVol", "VesselVol", "NonCALCMATXVol", "TotalPlaqueVolume"]:
        dec.append([OQ[v]["lab"], f'{R1[v]["delta"]:+.3f} → {R2[v]["delta"]:+.3f}',
                    f'{R1[v]["s2_rand"]:.4f} → {R2[v]["s2_rand"]:.4f}',
                    f'{R1[v]["wcv_rand_log"]:.1f} → {R2[v]["wcv_rand_log"]:.1f}'])
    table_slide("Key insight: bias moved, dispersion did not (v1 → v2)",
                ["Endpoint", "Δ bias (log)", "σ²_random", "wCV scanner-term %"], dec,
                colw=[3.0, 4.0, 3.6, 3.0])

    # 6 scanner-attributable — v2 canonical vs sub-segment
    def sacell(r):
        return f'{r["wcv_scan_log"]:.1f}  {"P" if r["scan_a_pass"] else "F"}/{"P" if r["scan_b_pass"] else "F"}'
    sa = [[OQ[v]["lab"], f'{OQ[v]["log"][0]}', sacell(R["v2c"][v]), sacell(R["v2s"][v])] for v in OQ]
    table_slide("Scanner-attributable variance vs OQ (v2): canonical vs sub-segment",
                ["Endpoint", "OQ", "canonical  wCV a/b", "sub-segment  wCV a/b"], sa,
                colw=[3.2, 1.6, 3.9, 3.9], size=13)

    # 6b Gate 4 systematic bias (log BA) vs OQ Table 6 — v2 canonical vs sub-segment
    br = []
    for v in PLAQUE:
        oqb = OQ[v]["bias_ci"]; rc, rs = R["v2c"][v], R["v2s"][v]
        br.append([OQ[v]["lab"], f'[{oqb[0]}, {oqb[1]}]',
                   f'{rc["delta"]:+.3f}  {"P" if rc["bias_overlap"] else "F"}',
                   f'{rs["delta"]:+.3f}  {"P" if rs["bias_overlap"] else "F"}'])
    table_slide("Gate 4 — systematic bias (log BA) vs OQ Table 6 (v2, plaque)",
                ["Endpoint", "OQ bias 95% CI", "canonical Δ (P/F)", "sub-segment Δ (P/F)"], br,
                colw=[3.0, 3.0, 3.4, 3.4], size=14)

    # 7 conclusions
    s = slide("Conclusions & future work")
    bullets(s, [
        ("Process outputs pass the OQ CI-overlap with the scanner term on BOTH vintages, canonical "
         "region — no sub-segment needed. Separating the systematic bias into Gate 4 is the key step.", 0, GREEN),
        ("Dispersion is vintage-stable; the 07-07 re-work increased systematic bias, not noise.", 0, None),
        ("Plaque: passes on untransformed scale (wide OQ CIs); CALC borderline on log scale — both reported.", 0, None),
        ("Scanner-attributable variance (Var(d)/2 − σ²_OQ) ≤ OQ / CI-includes-0 for most endpoints on "
         "canonical, and nearly all on the extent-matched sub-segment (regenerated on v2, N=24).", 0, None),
        ("After removing reader+repeat (OQ) and traced extent (sub-segment), the residual scanner "
         "variability is within the OQ envelope. Next: expand to N ≥ 30.", 0, BLUE),
    ], top=1.5, size=17)

    out = os.path.join(ROOT, "PCCT_Statistical_Methods_and_Results.pptx")
    prs.save(out); print("Wrote", out)


if __name__ == "__main__":
    build_docx()
    build_pptx()
